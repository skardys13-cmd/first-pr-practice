#!/usr/bin/env python3
"""
Raw extraction harness for the VC / PE / M&A / New Venture Financing
course knowledge project.

PURPOSE
    Convert an original course file (.pptx/.ppt, .pdf, .docx, .xlsx/.csv) into a
    complete, source-traceable raw text extraction under
    01_PRESENTATION_EXTRACTIONS/ (decks) or 04_DOCUMENT_NOTES/ (everything else).

    This produces the RAW layer only -- stage 1 of
    Extract -> Organize -> Synthesize -> Apply.
    Interpretation, synthesis and note-writing are done by Claude reading these
    extractions; this script never summarizes and never paraphrases.

HARD RULES
    * Originals are READ-ONLY. This script opens source files in 'rb' mode only
      and writes exclusively into derived output directories.
    * Nothing is invented. If a slide has no text, it is emitted as
      "[NO TEXT EXTRACTED]" rather than described.
    * Every extracted unit keeps its address (slide N / page N / sheet name /
      cell ref) so downstream notes can cite e.g. "SRC-P-004, Slide 11".

WHAT IT PULLS
    .pptx  slide text frames, speaker notes, tables, chart series/categories,
           SmartArt/diagram text, grouped-shape text, and an inventory of
           embedded images (optionally exported so they can be viewed).
    .pdf   per-page text plus detected tables.
    .docx  paragraphs (with heading levels) and tables.
    .xlsx  per-sheet cell values AND underlying formulas, plus defined names.
    .csv   rows as-is.
    .ppt/.doc/.xls  auto-converted to the modern format via LibreOffice into a
           scratch dir, then parsed (original untouched).

USAGE
    python3 extract_source.py --all
    python3 extract_source.py --file "03_SOURCE_DOCUMENTS/Lecture 04.pptx" --id SRC-P-004
    python3 extract_source.py --all --export-images   # also write slide images to disk
"""

import argparse
import csv
import os
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

# --------------------------------------------------------------------------
# Paths
# --------------------------------------------------------------------------

KB_ROOT = Path(__file__).resolve().parents[2]
SOURCE_DIR = KB_ROOT / "03_SOURCE_DOCUMENTS"
DECK_OUT = KB_ROOT / "01_PRESENTATION_EXTRACTIONS"
DOC_OUT = KB_ROOT / "04_DOCUMENT_NOTES"
IMAGE_OUT = KB_ROOT / "01_PRESENTATION_EXTRACTIONS" / "_extracted_images"

DECK_EXT = {".pptx", ".ppt"}
PDF_EXT = {".pdf"}
WORD_EXT = {".docx", ".doc"}
SHEET_EXT = {".xlsx", ".xlsm", ".xls"}
CSV_EXT = {".csv", ".tsv"}
TEXT_EXT = {".txt", ".md"}
SUPPORTED = DECK_EXT | PDF_EXT | WORD_EXT | SHEET_EXT | CSV_EXT | TEXT_EXT

LEGACY_CONVERT = {".ppt": "pptx", ".doc": "docx", ".xls": "xlsx"}


# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------

def slug(name: str) -> str:
    """Filesystem-safe stem."""
    return re.sub(r"[^A-Za-z0-9._-]+", "_", name).strip("_")


def convert_legacy(path: Path, scratch: Path) -> Path:
    """LibreOffice-convert a legacy Office file. Original is never modified."""
    target = LEGACY_CONVERT[path.suffix.lower()]
    subprocess.run(
        ["soffice", "--headless", "--convert-to", target, "--outdir", str(scratch), str(path)],
        check=True, capture_output=True, timeout=300,
    )
    out = scratch / (path.stem + "." + target)
    if not out.exists():
        raise RuntimeError(f"LibreOffice conversion produced no output for {path.name}")
    return out


def shape_text_blocks(shape, depth=0):
    """Recursively pull text from a shape, descending into groups."""
    blocks = []
    pad = "  " * depth

    if shape.shape_type is not None and str(shape.shape_type).startswith("GROUP"):
        for sub in shape.shapes:
            blocks.extend(shape_text_blocks(sub, depth + 1))
        return blocks

    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs).strip()
            if not line:
                continue
            bullet = "  " * para.level + "- " if para.level else ""
            blocks.append(f"{pad}{bullet}{line}")
    return blocks


# --------------------------------------------------------------------------
# Extractors
# --------------------------------------------------------------------------

def _extract_pptx(path: Path, src_id: str, export_images: bool) -> str:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(path))
    out = []
    total = len(prs.slides)
    out.append(f"# RAW EXTRACTION — {src_id}")
    out.append("")
    out.append(f"**Original file:** `{path.name}`  ")
    out.append(f"**Type:** PowerPoint deck  ")
    out.append(f"**Slide count:** {total}  ")
    try:
        w, h = prs.slide_width, prs.slide_height
        out.append(f"**Slide size:** {Emu(w).inches:.2f}in x {Emu(h).inches:.2f}in  ")
    except Exception:
        pass
    out.append("")
    out.append("> Machine extraction. Text is verbatim from the deck. "
               "Charts report their underlying plotted data. Images are inventoried "
               "(and exported when `--export-images` is used) so they can be read visually — "
               "an inventoried image has NOT been interpreted by this script.")
    out.append("")
    out.append("---")
    out.append("")

    img_count = 0
    for idx, slide in enumerate(prs.slides, start=1):
        out.append(f"## Slide {idx}")
        out.append("")

        try:
            layout = slide.slide_layout.name
            out.append(f"*Layout: {layout}*")
            out.append("")
        except Exception:
            pass

        body, tables, charts, images = [], [], [], []

        for shape in slide.shapes:
            # Tables
            if getattr(shape, "has_table", False):
                rows = []
                for r in shape.table.rows:
                    rows.append([c.text.strip().replace("\n", " ") for c in r.cells])
                tables.append(rows)
                continue

            # Charts
            if getattr(shape, "has_chart", False):
                ch = shape.chart
                info = [f"Chart type: {ch.chart_type}"]
                try:
                    cats = [str(c) for c in ch.plots[0].categories]
                    if cats:
                        info.append(f"Categories: {cats}")
                except Exception:
                    info.append("Categories: [NOT EXTRACTABLE]")
                try:
                    for s in ch.series:
                        vals = list(s.values)
                        info.append(f"Series '{s.name}': {vals}")
                except Exception:
                    info.append("Series: [NOT EXTRACTABLE]")
                try:
                    if ch.has_title:
                        info.insert(0, f"Chart title: {ch.chart_title.text_frame.text}")
                except Exception:
                    pass
                charts.append(info)
                continue

            # Pictures
            if str(shape.shape_type).startswith("PICTURE"):
                img_count += 1
                label = getattr(shape, "name", "picture")
                entry = f"{label}"
                if export_images:
                    try:
                        img = shape.image
                        IMAGE_OUT.mkdir(parents=True, exist_ok=True)
                        fn = IMAGE_OUT / f"{src_id}_slide{idx:03d}_img{img_count:03d}.{img.ext}"
                        fn.write_bytes(img.blob)
                        entry += f" -> exported `{fn.relative_to(KB_ROOT)}`"
                    except Exception as e:
                        entry += f" [EXPORT FAILED: {type(e).__name__}]"
                images.append(entry)
                continue

            body.extend(shape_text_blocks(shape))

        if body:
            out.append("### Text")
            out.append("")
            out.extend(body)
            out.append("")

        for t in tables:
            out.append("### Table")
            out.append("")
            if t:
                header = t[0]
                out.append("| " + " | ".join(header) + " |")
                out.append("|" + "---|" * len(header))
                for row in t[1:]:
                    out.append("| " + " | ".join(row) + " |")
            out.append("")

        for c in charts:
            out.append("### Chart data")
            out.append("")
            for line in c:
                out.append(f"- {line}")
            out.append("")

        if images:
            out.append("### Images on slide")
            out.append("")
            for i in images:
                out.append(f"- {i}")
            out.append("")
            out.append("> **[IMAGE NOT YET INTERPRETED]** — read the exported file(s) "
                       "visually before finalizing notes for this slide.")
            out.append("")

        # Speaker notes
        try:
            if slide.has_notes_slide:
                nt = slide.notes_slide.notes_text_frame.text.strip()
                if nt:
                    out.append("### Speaker notes")
                    out.append("")
                    out.append(nt)
                    out.append("")
        except Exception:
            pass

        if not (body or tables or charts or images):
            out.append("[NO TEXT EXTRACTED]")
            out.append("")

        out.append("---")
        out.append("")

    return "\n".join(out)


def _extract_pdf(path: Path, src_id: str) -> str:
    import pdfplumber

    out = [f"# RAW EXTRACTION — {src_id}", "",
           f"**Original file:** `{path.name}`  ", "**Type:** PDF  "]
    with pdfplumber.open(str(path)) as pdf:
        out.append(f"**Page count:** {len(pdf.pages)}  ")
        out.append("")
        out.append("---")
        out.append("")
        for i, page in enumerate(pdf.pages, start=1):
            out.append(f"## Page {i}")
            out.append("")
            text = (page.extract_text() or "").strip()
            out.append(text if text else "[NO TEXT EXTRACTED — may be a scanned image page]")
            out.append("")
            try:
                for tbl in page.extract_tables() or []:
                    out.append("### Table")
                    out.append("")
                    clean = [[(c or "").replace("\n", " ").strip() for c in row] for row in tbl]
                    if clean:
                        out.append("| " + " | ".join(clean[0]) + " |")
                        out.append("|" + "---|" * len(clean[0]))
                        for row in clean[1:]:
                            out.append("| " + " | ".join(row) + " |")
                    out.append("")
            except Exception:
                out.append("[TABLE EXTRACTION FAILED ON THIS PAGE]")
                out.append("")
            out.append("---")
            out.append("")
    return "\n".join(out)


def _extract_docx(path: Path, src_id: str) -> str:
    import docx

    d = docx.Document(str(path))
    out = [f"# RAW EXTRACTION — {src_id}", "",
           f"**Original file:** `{path.name}`  ", "**Type:** Word document  ", "",
           "---", ""]
    for para in d.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        style = (para.style.name or "").lower()
        if style.startswith("heading"):
            lvl = "".join(ch for ch in style if ch.isdigit()) or "1"
            out.append("#" * (int(lvl) + 1) + " " + t)
        elif "list" in style:
            out.append(f"- {t}")
        else:
            out.append(t)
        out.append("")

    for n, table in enumerate(d.tables, start=1):
        out.append(f"### Table {n}")
        out.append("")
        rows = [[c.text.strip().replace("\n", " ") for c in r.cells] for r in table.rows]
        if rows:
            out.append("| " + " | ".join(rows[0]) + " |")
            out.append("|" + "---|" * len(rows[0]))
            for row in rows[1:]:
                out.append("| " + " | ".join(row) + " |")
        out.append("")
    return "\n".join(out)


def _extract_xlsx(path: Path, src_id: str) -> str:
    import openpyxl

    out = [f"# RAW EXTRACTION — {src_id}", "",
           f"**Original file:** `{path.name}`  ", "**Type:** Spreadsheet / model  ", ""]
    out.append("> Both computed VALUES and underlying FORMULAS are captured. "
               "The formulas are the teaching content — they show the financial logic.")
    out.append("")

    wb_v = openpyxl.load_workbook(str(path), data_only=True, read_only=False)
    wb_f = openpyxl.load_workbook(str(path), data_only=False, read_only=False)

    out.append(f"**Sheets:** {wb_v.sheetnames}")
    out.append("")
    try:
        names = list(wb_f.defined_names.keys())
        if names:
            out.append(f"**Defined names:** {names}")
            out.append("")
    except Exception:
        pass
    out.append("---")
    out.append("")

    for name in wb_v.sheetnames:
        sv, sf = wb_v[name], wb_f[name]
        out.append(f"## Sheet: {name}")
        out.append("")
        out.append(f"*Dimensions: {sv.dimensions}*")
        out.append("")
        out.append("| Cell | Value | Formula |")
        out.append("|---|---|---|")
        empty = True
        for row_v, row_f in zip(sv.iter_rows(), sf.iter_rows()):
            for cv, cf in zip(row_v, row_f):
                if cv.value is None and cf.value is None:
                    continue
                empty = False
                val = str(cv.value).replace("|", "\\|") if cv.value is not None else ""
                fml = str(cf.value).replace("|", "\\|") if cf.value is not None else ""
                fml = fml if fml.startswith("=") else ""
                out.append(f"| {cv.coordinate} | {val} | {fml} |")
        if empty:
            out.append("| — | [SHEET EMPTY] | |")
        out.append("")
        out.append("---")
        out.append("")
    return "\n".join(out)


def _extract_text(path: Path, src_id: str) -> str:
    """
    Pasted lecture text. Preserved VERBATIM -- never reflowed, trimmed, or
    reworded. Slide/section boundaries are detected and indexed where the paste
    contains them, so downstream notes can still cite a slide number.
    """
    raw = path.read_text(encoding="utf-8", errors="replace")
    lines = raw.splitlines()

    # Detect slide-ish delimiters so citations stay addressable.
    marker = re.compile(
        r"^\s*(?:#+\s*)?(?:slide|lecture|section|part|chapter|module)\s*[#:\-]?\s*(\d+)\b",
        re.IGNORECASE,
    )
    hits = [(i + 1, m.group(0).strip()) for i, ln in enumerate(lines)
            if (m := marker.match(ln))]

    out = [f"# RAW EXTRACTION — {src_id}", "",
           f"**Original file:** `{path.name}`  ",
           "**Type:** Pasted lecture text  ",
           f"**Lines:** {len(lines)}  ",
           f"**Detected slide/section markers:** {len(hits)}  ", ""]
    out.append("> Verbatim capture of pasted course content. Nothing below is "
               "summarized, reworded, or reordered. If the paste came from a deck, "
               "formatting that carried meaning in the original (chart data, table "
               "layout, images) may be absent — check `MARKERS` below for how "
               "granular the citations can be.")
    out.append("")

    if hits:
        out.append("## Detected markers")
        out.append("")
        for ln_no, label in hits:
            out.append(f"- line {ln_no}: `{label}`")
        out.append("")
    else:
        out.append("> **[NO SLIDE MARKERS DETECTED]** — citations for this source "
                   "must reference line numbers rather than slide numbers, unless "
                   "slide structure can be inferred from the content itself.")
        out.append("")

    out.append("---")
    out.append("")
    out.append("## Verbatim content")
    out.append("")
    out.append("```text")
    out.extend(lines)
    out.append("```")
    return "\n".join(out)


def _extract_csv(path: Path, src_id: str) -> str:
    out = [f"# RAW EXTRACTION — {src_id}", "",
           f"**Original file:** `{path.name}`  ", "**Type:** Delimited data  ", "",
           "---", ""]
    delim = "\t" if path.suffix.lower() == ".tsv" else ","
    with path.open(newline="", encoding="utf-8", errors="replace") as fh:
        rows = list(csv.reader(fh, delimiter=delim))
    if not rows:
        out.append("[FILE EMPTY]")
        return "\n".join(out)
    out.append("| " + " | ".join(rows[0]) + " |")
    out.append("|" + "---|" * len(rows[0]))
    for row in rows[1:]:
        out.append("| " + " | ".join(row) + " |")
    return "\n".join(out)


# --------------------------------------------------------------------------
# Driver
# --------------------------------------------------------------------------

def process(path: Path, src_id: str, export_images: bool, scratch: Path) -> Path:
    ext = path.suffix.lower()
    work = path

    if ext in LEGACY_CONVERT:
        print(f"    converting legacy {ext} via LibreOffice ...")
        work = convert_legacy(path, scratch)
        ext = work.suffix.lower()

    if ext == ".pptx":
        body, dest_dir = _extract_pptx(work, src_id, export_images), DECK_OUT
    elif ext == ".pdf":
        body, dest_dir = _extract_pdf(work, src_id), DOC_OUT
    elif ext == ".docx":
        body, dest_dir = _extract_docx(work, src_id), DOC_OUT
    elif ext in {".xlsx", ".xlsm"}:
        body, dest_dir = _extract_xlsx(work, src_id), DOC_OUT
    elif ext in CSV_EXT:
        body, dest_dir = _extract_csv(work, src_id), DOC_OUT
    elif ext in TEXT_EXT:
        # Pasted lectures are presentation sources; other pasted docs are not.
        body = _extract_text(work, src_id)
        dest_dir = DECK_OUT if src_id.startswith("SRC-P") else DOC_OUT
    else:
        raise ValueError(f"unsupported extension: {ext}")

    dest_dir.mkdir(parents=True, exist_ok=True)
    dest = dest_dir / f"{src_id}_{slug(path.stem)}_RAW.md"
    dest.write_text(body, encoding="utf-8")
    return dest


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--all", action="store_true",
                    help=f"process every supported file in {SOURCE_DIR.name}/")
    ap.add_argument("--file", help="process a single file")
    ap.add_argument("--id", help="source ID for --file (e.g. SRC-P-004)")
    ap.add_argument("--export-images", action="store_true",
                    help="write embedded deck images to disk for visual reading")
    args = ap.parse_args()

    if args.file:
        p = Path(args.file)
        if not p.is_absolute():
            p = (KB_ROOT / args.file) if (KB_ROOT / args.file).exists() else Path.cwd() / args.file
        targets = [(p, args.id or f"SRC-X-{slug(p.stem)[:20]}")]
    elif args.all:
        if not SOURCE_DIR.exists():
            print(f"ERROR: {SOURCE_DIR} does not exist.", file=sys.stderr)
            return 1
        # Skip this folder's own control/readme files -- they are not course sources.
        files = sorted(f for f in SOURCE_DIR.rglob("*")
                       if f.is_file()
                       and f.suffix.lower() in SUPPORTED
                       and not f.name.startswith((".", "_"))
                       and not f.stem.upper().startswith("README"))
        if not files:
            print(f"No supported source files found in {SOURCE_DIR}.")
            print("Drop the course decks/PDFs/docs there, then re-run.")
            return 0
        targets = []
        deck_n = doc_n = 0
        for f in files:
            if f.suffix.lower() in DECK_EXT | TEXT_EXT:
                deck_n += 1
                targets.append((f, f"SRC-P-{deck_n:03d}"))
            else:
                doc_n += 1
                targets.append((f, f"SRC-D-{doc_n:03d}"))
    else:
        ap.print_help()
        return 1

    scratch = Path(tempfile.mkdtemp(prefix="vckb_"))
    ok = fail = 0
    try:
        for path, src_id in targets:
            print(f"[{src_id}] {path.name}")
            try:
                dest = process(path, src_id, args.export_images, scratch)
                print(f"    -> {dest.relative_to(KB_ROOT)}")
                ok += 1
            except Exception as e:
                print(f"    !! FAILED: {type(e).__name__}: {e}", file=sys.stderr)
                fail += 1
    finally:
        shutil.rmtree(scratch, ignore_errors=True)

    print(f"\nDone. {ok} extracted, {fail} failed.")
    print("NOTE: assigned IDs are provisional — set final IDs in FILE_MANIFEST.md "
          "based on actual course order, not filename order.")
    return 0 if fail == 0 else 2


if __name__ == "__main__":
    sys.exit(main())
